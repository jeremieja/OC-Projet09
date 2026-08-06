"""
Remplace la diapositive unique « Les sources qui fondent ces choix » par quatre
diapositives, une par référence, dans le diaporama de soutenance existant.

Le fichier est modifié SUR PLACE, après création d'une sauvegarde horodatée.
La charte est préservée : on réutilise la mise en page « Résumé 2 » du diaporama
et la géométrie exacte des zones de texte de la diapositive d'origine.

Usage : python scripts/eclater_slide_sources.py [chemin_du_pptx]
"""
import shutil
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

DEFAUT = Path("C:/Users/JEREMIE/Desktop/Jambon_Jeremie_6_presentation_082026.pptx")

# Position de la diapositive à remplacer (1-indexée, telle qu'affichée dans PowerPoint)
SLIDE_A_REMPLACER = 8
SECTION = "1 – Plan prévisionnel"

INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x00, 0x72, 0xB2)
TIP_BG = RGBColor(0xEF, 0xF4, 0xFA)

# Géométrie relevée sur la diapositive d'origine, pour un rendu identique.
ZONE_CONTENU = (0.62, 1.45, 12.10, 3.90)
ZONE_ENCART = (8.82, 5.75, 3.90, 1.35)

SOURCES = [
    {
        "sous_titre": "Source 1 — L'article de recherche fondateur",
        "puces": [
            "Tunstall et al. (2022), « Efficient Few-Shot Learning Without Prompts », "
            "arXiv:2209.11055.",
            "Publié conjointement par Hugging Face, Intel Labs et le UKP Lab de "
            "l'université technique de Darmstadt.",
            "Ce qu'il apporte : la méthode SetFit en deux étapes — affinage par "
            "comparaison de paires, puis classifieur léger.",
            ("Il répond à une limite des approches antérieures, qui dépendaient de "
             "consignes textuelles rédigées à la main et de modèles géants.",),
            "Le résultat de référence : 8 exemples par classe rivalisent avec un "
            "affinage classique sur 3 000 exemples.",
        ],
        "encart": "💡 C'est la référence obligatoire du projet : elle décrit la méthode "
                  "et fournit les résultats auxquels je compare les miens.",
    },
    {
        "sous_titre": "Source 2 — La mise en pratique officielle",
        "puces": [
            "Hugging Face (2022), article de blog « SetFit », huggingface.co/blog/setfit.",
            "Rédigé par les auteurs eux-mêmes : c'est la traduction opérationnelle de "
            "l'article, avec du code exécutable.",
            "Ce qu'il apporte : la librairie setfit (licence libre) et la structure "
            "d'entraînement que j'ai réutilisée.",
            "Il détaille aussi le banc d'essai RAFT : SetFit, avec 355 millions de "
            "paramètres, obtient 71,3 % contre 62,7 % pour GPT-3 et ses 175 milliards.",
        ],
        "encart": "💡 Transparence : c'est de cette source que vient la structure "
                  "d'entraînement de mon code. Le corpus, le protocole, les stratégies 4 "
                  "et 5 et le tableau de bord sont originaux.",
    },
    {
        "sous_titre": "Source 3 — Le banc d'essai qui justifie le socle",
        "puces": [
            "Ciancone et al. (2024), « Extending the Massive Text Embedding Benchmark "
            "to French », arXiv:2405.20468.",
            "Premier banc d'essai massif de représentations de phrases en français.",
            ("Une cinquantaine de modèles comparés, sur 8 familles de tâches et "
             "18 jeux de données.",),
            "Sa conclusion : aucun modèle ne gagne partout, mais les bons modèles "
            "multilingues entraînés sur la similarité de phrases rivalisent avec les "
            "modèles français spécialisés.",
        ],
        "encart": "💡 C'est ce qui fonde mon choix d'un socle multilingue plutôt que d'un "
                  "modèle exclusivement français.",
    },
    {
        "sous_titre": "Les sources complémentaires",
        "puces": [
            "Reimers & Gurevych (2020), arXiv:2004.09813 — la méthode de distillation "
            "qui a produit le socle multilingue que j'utilise.",
            "Warner et al. (2024), arXiv:2412.13663 — ModernBERT, première modernisation "
            "de l'architecture BERT depuis l'originale.",
            "Marone et al. (2025) — mmBERT, la version multilingue de ModernBERT, "
            "entraînée sur plus de 1 800 langues.",
            "Martin et al. (2020), arXiv:1911.03894 — CamemBERT, le modèle français "
            "utilisé comme stratégie 2.",
        ],
        "encart": "💡 Critères retenus pour qualifier un algorithme de « récent » : moins "
                  "de cinq ans, publié sur un support reconnu, et disponible dans une "
                  "librairie libre et maintenue.",
    },
]


def titre_section(slide, texte):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = texte
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(20)
            return


def bloc_texte(slide, sous_titre, puces):
    x, y, w, h = ZONE_CONTENU
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = sous_titre
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    p.space_after = Pt(12)

    for item in puces:
        niveau = 1 if isinstance(item, tuple) else 0
        txt = item[0] if isinstance(item, tuple) else item
        par = tf.add_paragraph()
        par.level = niveau
        run = par.add_run()
        run.text = ("• " if niveau == 0 else "– ") + txt
        run.font.size = Pt(14) if niveau == 0 else Pt(12.5)
        run.font.color.rgb = INK if niveau == 0 else INK_SOFT
        par.space_after = Pt(8)


def encart(slide, texte):
    x, y, w, h = ZONE_ENCART
    forme = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    forme.fill.solid()
    forme.fill.fore_color.rgb = TIP_BG
    forme.line.color.rgb = TIP_BG
    forme.shadow.inherit = False
    tf = forme.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = texte
    r.font.size = Pt(11)
    r.font.color.rgb = INK


def main():
    chemin = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAUT

    if not chemin.exists():
        sys.exit(f"Introuvable : {chemin}")
    verrou = chemin.parent / f"~${chemin.name}"
    if verrou.exists():
        sys.exit("Le diaporama est ouvert dans PowerPoint. Ferme-le (en enregistrant "
                 "tes modifications) puis relance ce script.")

    sauvegarde = chemin.with_name(
        f"{chemin.stem}_sauvegarde_{datetime.now():%Y%m%d_%H%M}.pptx")
    shutil.copy2(chemin, sauvegarde)
    print(f"Sauvegarde : {sauvegarde.name}")

    prs = Presentation(str(chemin))
    lst = prs.slides._sldIdLst
    total_avant = len(lst)
    if not 1 <= SLIDE_A_REMPLACER <= total_avant:
        sys.exit(f"La diapositive {SLIDE_A_REMPLACER} n'existe pas ({total_avant} au total).")

    # Mise en page réutilisée : celle de la diapositive remplacée, pour conserver la charte.
    layout = list(prs.slides)[SLIDE_A_REMPLACER - 1].slide_layout

    # 1. Créer les nouvelles diapositives (elles s'ajoutent en fin de diaporama).
    for src in SOURCES:
        s = prs.slides.add_slide(layout)
        titre_section(s, SECTION)
        bloc_texte(s, src["sous_titre"], src["puces"])
        encart(s, src["encart"])

    # 2. Supprimer l'ancienne diapositive : on retire l'entrée de la liste ET la
    #    relation qui y mène, sinon elle resterait physiquement dans le fichier.
    R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    ancienne = lst[SLIDE_A_REMPLACER - 1]
    prs.part.drop_rel(ancienne.get(R_ID))
    lst.remove(ancienne)

    # 3. Déplacer les quatre nouvelles diapositives à la place de l'ancienne.
    nouvelles = list(lst)[-len(SOURCES):]
    for n in nouvelles:
        lst.remove(n)
    for decalage, n in enumerate(nouvelles):
        lst.insert(SLIDE_A_REMPLACER - 1 + decalage, n)

    prs.save(str(chemin))
    print(f"Diapositives : {total_avant} -> {len(prs.slides._sldIdLst)}")
    print(f"Modifié : {chemin.name}")


if __name__ == "__main__":
    main()
