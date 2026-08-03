"""
Génère le support de soutenance (.pptx) en réutilisant la charte graphique
du modèle fourni : on ouvre le fichier existant, on retire ses diapositives
mais on conserve le masque, les mises en page, les polices et les couleurs.

Structure (28 diapositives, limite de 30 respectée) :
  Couverture / Sommaire / Introduction
  1 - Élaboration du modèle
  2 - Résultats comparés
  3 - Dashboard et mise en production
  Conclusion / Merci

Les emplacements de captures d'écran du dashboard sont laissés vides et légendés :
il suffit de coller l'image par-dessus dans PowerPoint.

Sortie : Jambon_Jeremie_9_presentation.pptx
Usage  : python scripts/generate_presentation.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).parents[1]
TEMPLATE = ROOT / "Jambon_Jeremie_6_presentation_082026.pptx"
OUT = ROOT / "Jambon_Jeremie_9_presentation.pptx"
FIG = ROOT / "results" / "figures"
INTERP = ROOT / "results" / "interpretability"

# ── Informations de couverture (à ajuster si besoin) ─────────────────────────
TITRE = "Classification automatique des mails d'un club sportif"
ETUDIANT = "Jérémie JAMBON"
MENTOR = "[à compléter]"
DATE_SOUTENANCE = "[à compléter]"

# ── Repères graphiques ───────────────────────────────────────────────────────
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x00, 0x72, 0xB2)
TIP_BG = RGBColor(0xEF, 0xF4, 0xFA)
FRAME = RGBColor(0xBB, 0xBB, 0xBB)

# Mises en page du modèle
L_COVER, L_SOMMAIRE, L_CONTENU, L_FIN = 0, 1, 4, 12

# Intitulés calqués sur le plan de soutenance imposé, pour que le jury retrouve
# immédiatement les trois temps attendus et leur minutage.
SECTIONS = {
    "intro": "Introduction",
    "s1": "1 – Plan prévisionnel",
    "s2": "2 – Démarche mise en œuvre",
    "s3": "3 – Démonstration du dashboard",
    "ccl": "Conclusion",
}


# ── Utilitaires de construction ──────────────────────────────────────────────
def vider(prs):
    """
    Retire les diapositives du modèle en conservant masque, mises en page et thème.

    Il ne suffit pas de retirer l'entrée de la liste : la diapositive resterait
    physiquement dans le fichier (avec ses images). On supprime donc aussi la
    relation qui y mène, ce qui rend la partie orpheline et l'exclut de la
    sauvegarde.
    """
    R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        prs.part.drop_rel(sld.get(R_ID))
        lst.remove(sld)


def _set_ph(slide, idx, texte):
    """Renseigne un placeholder s'il existe (index de position dans le layout)."""
    try:
        ph = slide.placeholders[idx]
    except (KeyError, IndexError):
        return None
    ph.text = texte
    return ph


def titre_section(slide, texte):
    """Le titre de la mise en page porte le nom de la section courante."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = texte
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(20)
            return


def bloc_texte(slide, sous_titre, puces, x=0.62, y=1.45, w=6.3, h=4.3):
    """Zone de contenu : un sous-titre puis des puces de niveau 1 ou 2."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = sous_titre
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    p.space_after = Pt(10)

    for item in puces:
        niveau = 1 if isinstance(item, tuple) else 0
        txt = item[0] if isinstance(item, tuple) else item
        par = tf.add_paragraph()
        par.level = niveau
        run = par.add_run()
        run.text = ("• " if niveau == 0 else "– ") + txt
        run.font.size = Pt(14) if niveau == 0 else Pt(12.5)
        run.font.color.rgb = INK if niveau == 0 else INK_SOFT
        par.space_after = Pt(6)
    return tb


def encart(slide, texte, y=6.05, x=0.62, w=12.1, h=1.05):
    """Encart pédagogique repris du modèle (fond clair, texte court)."""
    forme = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    forme.fill.solid()
    forme.fill.fore_color.rgb = TIP_BG
    forme.line.color.rgb = TIP_BG
    forme.shadow.inherit = False

    tf = forme.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = texte
    r.font.size = Pt(12)
    r.font.color.rgb = INK
    return forme


def image(slide, chemin, x, y, w):
    if Path(chemin).exists():
        slide.shapes.add_picture(str(chemin), Inches(x), Inches(y), width=Inches(w))


def emplacement_capture(slide, legende, x=7.05, y=1.5, w=5.7, h=4.2):
    """Cadre vide destiné à recevoir une capture d'écran du dashboard."""
    forme = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    forme.fill.background()
    forme.line.color.rgb = FRAME
    forme.line.width = Pt(1.25)
    forme.line.dash_style = 2  # pointillés
    forme.shadow.inherit = False
    tf = forme.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"[ Capture d'écran à insérer ]\n{legende}"
    r.font.size = Pt(12)
    r.font.color.rgb = INK_SOFT
    return forme


def nouvelle(prs, layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


# ── Construction du support ──────────────────────────────────────────────────
def build():
    prs = Presentation(str(TEMPLATE))
    vider(prs)

    # 1 — Couverture
    s = nouvelle(prs, L_COVER)
    titre_section(s, TITRE)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(5.3), Inches(7.0), Inches(1.2))
    tf = tb.text_frame
    for i, ligne in enumerate([f"Étudiant : {ETUDIANT}",
                               f"Mentor : {MENTOR}",
                               f"Date de soutenance : {DATE_SOUTENANCE}"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ligne
        r.font.size = Pt(14); r.font.color.rgb = INK_SOFT

    # 2 — Sommaire
    s = nouvelle(prs, L_SOMMAIRE)
    titre_section(s, "Sommaire")
    entrees = ["Introduction",
               "1 – Plan prévisionnel  (5 min)",
               "2 – Démarche mise en œuvre  (10 min)",
               "3 – Démonstration du dashboard  (5 min)",
               "Conclusion"]
    # On cible le placeholder de contenu par son type (et non par un index codé en
    # dur, qui varie d'une mise en page à l'autre).
    corps = next((ph for ph in s.placeholders
                  if ph.placeholder_format.idx not in (0,)
                  and ph.placeholder_format.type not in (PP_PLACEHOLDER.SLIDE_NUMBER,
                                                         PP_PLACEHOLDER.DATE)), None)
    if corps is not None:
        tf = corps.text_frame
        tf.clear()
        for i, e in enumerate(entrees):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = e
            r.font.size = Pt(20); r.font.color.rgb = INK
            p.space_after = Pt(12)
    else:
        bloc_texte(s, "", entrees, y=1.8, w=9.0)

    # 3 — Introduction : le problème métier
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["intro"])
    bloc_texte(s, "Le contexte", [
        "Un logiciel destiné aux clubs sportifs structurés (football, handball, volley…).",
        "Sa brique centrale : trier et prioriser automatiquement les mails entrants.",
        "Huit catégories métier : inscription, sponsor, arbitrage, parent, fédération, "
        "logistique de match, indemnités, administratif.",
        "La difficulté : un club qui arrive sur le logiciel n'a aucun mail déjà classé.",
    ], w=12.1)
    encart(s, "💡 Problématique : comment obtenir un classificateur fiable quand on ne dispose "
              "que de quelques exemples étiquetés ? C'est le cadre du « few-shot learning ».")

    # 4 — §1 Les données
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s1"])
    bloc_texte(s, "Deux jeux de données complémentaires", [
        "Cas d'usage métier : 1 800 mails de clubs sportifs en français, 8 catégories.",
        ("225 mails par catégorie, corpus parfaitement équilibré.",),
        ("Découpage stratifié : 1 440 pour l'entraînement, 360 pour le test.",),
        "Validation scientifique : 20 Newsgroups, 8 classes, ~7 200 documents.",
        ("Benchmark public reconnu : il situe nos résultats par rapport à la littérature.",),
    ], w=12.1, h=2.4)
    image(s, FIG / "fig_eda_emails.png", x=0.75, y=3.7, w=11.8)

    # 6 — §1 Génération du dataset
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s1"])
    bloc_texte(s, "Construction du jeu de données métier", [
        "Aucun corpus public de mails de clubs sportifs n'existe : il a fallu le créer.",
        "Génération par grand modèle de langue, à partir de consignes structurées.",
        "Variations contrôlées pour éviter des mails stéréotypés :",
        ("longueur, ton (formel, direct, urgent, hésitant),",),
        ("profil de l'expéditeur (parent, entraîneur, dirigeant, fédération…),",),
        ("sport concerné (10 disciplines différentes).",),
        "Sauvegarde catégorie par catégorie, pour reprendre en cas d'interruption.",
    ], w=12.1)
    encart(s, "💡 Un jeu de données synthétique permet de travailler sans exposer de données "
              "personnelles réelles — un atout au regard du RGPD.")

    # 7 — §1 Les cinq stratégies
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s1"])
    bloc_texte(s, "Cinq stratégies, une logique produit", [
        "1. TF-IDF + régression logistique — la référence historique, sans carte graphique.",
        "2. CamemBERT affiné — modèle de langue français, approche « pilotée par la donnée ».",
        "3. SetFit — l'algorithme récent étudié, conçu pour les faibles volumes.",
        "4. Ministral 8B via API — le modèle génératif utilisé directement comme classifieur.",
        "5. Système hybride — SetFit en première ligne, le modèle génératif en renfort.",
    ], w=12.1)
    encart(s, "💡 L'objectif n'est pas seulement de comparer des scores, mais de répondre à une "
              "question de produit : quelle stratégie déployer selon la maturité du club ?")

    # §1 — Justification du modèle retenu
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s1"])
    bloc_texte(s, "Pourquoi SetFit ? Ce que dit la littérature", [
        "Article fondateur : Tunstall et al. (2022), Hugging Face, Intel Labs et UKP Lab.",
        "Avec 8 exemples par classe, SetFit rivalise avec un modèle affiné sur 3 000 exemples.",
        "Sur le banc d'essai RAFT : 71,3 % contre 62,7 % pour GPT-3, un modèle 30 fois plus gros.",
        "Entraînement de l'ordre de la trentaine de secondes, pour quelques centimes.",
        "Aucune dépendance externe : le modèle tourne en local et l'inférence est gratuite.",
    ], w=12.1)
    encart(s, "💡 Trois variantes de socle ont été comparées : un modèle multilingue de "
              "référence, ModernBERT (2024) et mmBERT (2025), pour mesurer l'apport "
              "d'architectures plus récentes.")

    # §1 — Les sources bibliographiques (exigé par le plan de soutenance)
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s1"])
    bloc_texte(s, "Les sources qui fondent ces choix", [
        "Tunstall et al. (2022), « Efficient Few-Shot Learning Without Prompts », "
        "arXiv:2209.11055 — article de recherche fondateur de SetFit.",
        "Hugging Face (2022), article de blog officiel sur SetFit — résultats illustrés "
        "et code de référence, réutilisé pour structurer l'entraînement.",
        "Ciancone et al. (2024), « Extending MTEB to French », arXiv:2405.20468 — "
        "banc d'essai qui a guidé le choix du socle d'embeddings francophone.",
        "Compléments : Reimers & Gurevych (2020) pour le socle multilingue, "
        "Warner et al. (2024) pour ModernBERT, Marone et al. (2025) pour mmBERT.",
    ], w=12.1)
    encart(s, "💡 Les modèles retenus datent tous de moins de cinq ans, sont publiés sur "
              "des supports reconnus et disponibles dans des bibliothèques libres et "
              "maintenues — critères exigés pour un algorithme « récent ».")

    # 9 — §2 SetFit : principe (nouveau concept)
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "SetFit : apprendre à partir de quelques exemples", [
        "Étape 1 — Apprentissage par comparaison :",
        ("on forme des paires de mails : même catégorie, ou catégories différentes ;",),
        ("le modèle apprend à rapprocher les premières et à éloigner les secondes.",),
        "Étape 2 — Un classifieur léger (régression logistique) est entraîné sur les "
        "représentations ainsi obtenues.",
        "L'astuce : 128 mails produisent plusieurs milliers de paires d'entraînement.",
    ], w=12.1)
    encart(s, "💡 Plutôt que d'apprendre « ce mail est une inscription », SetFit apprend d'abord "
              "« ces deux mails se ressemblent ». C'est ce qui démultiplie le signal disponible.")

    # §2 — Le système hybride (nouveau concept)
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Une cinquième voie : le routage par confiance", [
        "SetFit traite localement la grande majorité des mails, gratuitement.",
        "Chaque prédiction s'accompagne d'un score de confiance.",
        "Si la confiance passe sous un seuil τ, le mail est transmis au modèle génératif.",
        "On ne paie l'API que sur les cas réellement ambigus.",
        "Le seuil τ devient un curseur produit entre qualité, coût et latence.",
    ], w=12.1)
    encart(s, "💡 C'est l'architecture qu'un éditeur logiciel déploierait réellement : un modèle "
              "local rapide en première ligne, une escalade ciblée vers un service payant.")

    # §2 — Protocole expérimental
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Protocole d'évaluation", [
        "Cinq régimes d'entraînement : 8, 16, 32, 64 exemples par catégorie, puis données complètes.",
        "Cinq tirages différents par régime, pour mesurer la stabilité des résultats.",
        "Environ 380 expériences au total.",
        "Métrique principale : le F1 macro.",
        ("moyenne des F1 calculés catégorie par catégorie ;",),
        ("une petite catégorie pèse autant qu'une grande, ce qui évite de masquer "
         "les erreurs sur les cas rares.",),
    ], w=12.1)
    encart(s, "💡 Faire varier le nombre d'exemples permet de tracer une courbe d'apprentissage : "
              "on observe le comportement de chaque modèle selon la maturité du club.")

    # 12 — §2 Courbes emails
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Sur le cas d'usage métier", [
        "SetFit atteint 0,95 dès 8 exemples par catégorie.",
        "CamemBERT s'effondre à 0,10 : trop de paramètres pour si peu de données.",
        "Les écarts se resserrent à mesure que le volume augmente.",
    ], w=5.6, h=2.2)
    image(s, FIG / "fig_courbes_emails.png", x=6.0, y=1.35, w=7.0)
    encart(s, "💡 Lecture : plus une courbe est haute à gauche, plus le modèle est efficace "
              "lorsque les données étiquetées sont rares.")

    # 13 — §2 Courbes newsgroups
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Vérification sur un banc d'essai public", [
        "Même hiérarchie que sur les mails : SetFit domine à faibles volumes.",
        "Les textes sont ici plus longs et plus bruités qu'un mail.",
        "Le constat n'est donc pas un artefact de notre jeu de données.",
    ], w=5.6, h=2.2)
    image(s, FIG / "fig_courbes_newsgroups.png", x=6.0, y=1.35, w=7.0)

    # 14 — §2 Full data
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Avec l'intégralité des données", [
        "Sur les mails, les quatre approches se rejoignent autour de 0,99.",
        "Le vocabulaire métier est très discriminant : « vestiaires », « licence », « partenariat »…",
        "La méthode historique suffit donc lorsque l'historique est fourni.",
    ], w=12.1, h=1.9)
    image(s, FIG / "fig_fulldata.png", x=1.0, y=3.1, w=11.3)

    # 15 — §2 Ministral
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Le modèle génératif comme classifieur", [
        "Aucun entraînement : on décrit la tâche en français dans la consigne.",
        "Sans aucun exemple : 0,79 sur les mails.",
        "Avec 8 exemples glissés dans la consigne : 0,90.",
        "Latence d'environ 400 à 600 millisecondes par mail, et un coût par appel.",
        "Intérêt : opérationnel immédiatement, sans la moindre donnée étiquetée.",
    ], w=12.1)
    encart(s, "💡 Limite observée : la confiance annoncée par le modèle génératif est "
              "auto-déclarée et peu fiable. C'est pourquoi le système hybride s'appuie sur "
              "la confiance de SetFit, et non sur la sienne.")

    # 16 — §2 Hybride
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Le meilleur compromis qualité / coût", [
        "Au seuil τ = 0,8 : F1 de 0,997, le meilleur score de l'étude.",
        "Seuls 1,1 % des mails sont transmis au modèle payant.",
        "Le gain de qualité est donc obtenu pour un coût quasi nul.",
    ], w=12.1, h=1.7)
    image(s, FIG / "fig_hybride.png", x=1.0, y=2.9, w=11.3)

    # 17 — §2 Le bug révélé
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Un résultat aberrant, et ce qu'il a révélé", [
        "Premier constat : le modèle génératif plafonnait à 0,18 sur le banc d'essai public — "
        "à peine mieux que le hasard.",
        "Vérification des réponses brutes : le modèle classait en réalité correctement.",
        "En remontant la chaîne : les étiquettes étaient décalées par rapport aux textes.",
        ("scikit-learn trie les catégories par ordre alphabétique, indépendamment de "
         "l'ordre demandé.",),
        "Après correction : 0,82. Les modèles entraînés, eux, masquaient l'anomalie.",
    ], w=12.1)
    encart(s, "💡 Un modèle non entraîné juge sur le contenu réel : il sert de détecteur "
              "d'incohérence de la vérité-terrain. C'est le contrôle qualité qui a sauvé "
              "la fiabilité de l'étude.")

    # 18 — §2 Interprétabilité globale
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Sur quoi le modèle fonde-t-il ses décisions ?", [
        "La régression logistique est directement lisible : chaque mot porte un poids.",
        "Les mots décisifs correspondent au métier, ce qui explique la robustesse observée.",
    ], w=12.1, h=1.5)
    image(s, INTERP / "global_tfidf.png", x=0.9, y=2.5, w=11.5)

    # 19 — §2 Interprétabilité locale
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Expliquer une prédiction, mail par mail", [
        "Exemple : un mail classé « logistique de match ».",
        "Régression logistique — mots décisifs : vestiaires, terrain, match, samedi.",
        "SetFit, expliqué par la méthode LIME : vestiaires, samedi, match, terrain.",
        "Les deux familles de modèles s'appuient donc sur les mêmes indices.",
    ], w=12.1)
    encart(s, "💡 LIME explique un modèle opaque en observant comment sa prédiction varie "
              "lorsqu'on retire des mots. La convergence avec un modèle transparent renforce "
              "la confiance dans les deux.")

    # 20 — §2 Conclusion intermédiaire
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s2"])
    bloc_texte(s, "Ce que montrent les résultats", [
        "SetFit est le plus efficace lorsque les données étiquetées sont rares.",
        "CamemBERT n'exprime son potentiel qu'avec un historique conséquent.",
        "La méthode historique reste imbattable en rapport performance / simplicité "
        "lorsque les données sont abondantes.",
        "Le modèle génératif permet un démarrage immédiat, sans aucune donnée.",
        "Le système hybride offre le meilleur compromis global.",
    ], w=12.1)

    # 21 — §3 Dashboard : objectif
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s3"])
    bloc_texte(s, "Un tableau de bord pour rendre le travail tangible", [
        "Trois volets : exploration des données, prédiction en direct, performances.",
        "Modèle servi en ligne : TF-IDF + régression logistique.",
        ("quelques mégaoctets, pas de carte graphique, réponse immédiate ;",),
        ("performance équivalente aux modèles lourds sur données complètes.",),
        "Comparaison optionnelle avec le modèle génératif, via son API.",
    ], w=6.2)
    emplacement_capture(s, "Vue générale du tableau de bord")
    encart(s, "💡 Choix d'ingénierie assumé : déployer le modèle le plus léger à performance "
              "équivalente, plutôt que le plus sophistiqué.")

    # 22 — §3 Exploration
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s3"])
    bloc_texte(s, "Volet 1 — Exploration des données", [
        "Statistiques descriptives et répartition des catégories.",
        "Distribution de la longueur des messages.",
        "Mots les plus fréquents, filtrables par catégorie.",
        "Nuage de mots généré à la volée.",
        "Graphiques interactifs : survol, zoom, sélection.",
    ], w=6.2)
    emplacement_capture(s, "Onglet « Exploration des données »")

    # 23 — §3 Prédiction
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s3"])
    bloc_texte(s, "Volet 2 — Prédiction en direct", [
        "Saisie libre d'un message, ou tirage aléatoire dans le jeu de données.",
        "Catégorie prédite, niveau de confiance et probabilités détaillées.",
        "Lorsque le mail provient du jeu de test, la bonne réponse est affichée : "
        "on vérifie immédiatement si le modèle a vu juste.",
        "Comparaison possible avec le modèle génératif.",
    ], w=6.2)
    emplacement_capture(s, "Onglet « Prédiction en direct »")

    # 24 — §3 Performances
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s3"])
    bloc_texte(s, "Volet 3 — Performances des modèles", [
        "Courbes d'apprentissage comparées, avec dispersion entre tirages.",
        "Tableau récapitulatif sur données complètes.",
        "Sélection du jeu de données à afficher.",
    ], w=6.2)
    emplacement_capture(s, "Onglet « Performances des modèles »")

    # 25 — §3 Déploiement et accessibilité
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["s3"])
    bloc_texte(s, "Mise en production et accessibilité", [
        "Code versionné sur GitHub, application hébergée sur Streamlit Cloud.",
        "Dépendances allégées pour le déploiement : le nécessaire, rien de plus.",
        "Clé d'API gérée par le mécanisme de secrets, jamais dans le code.",
        "Accessibilité (référentiel WCAG) :",
        ("palette conçue pour les déficiences de vision des couleurs ;",),
        ("aucune information portée par la couleur seule (libellés, valeurs, formes) ;",),
        ("navigation au clavier et descriptions textuelles des visuels.",),
    ], w=12.1)

    # 26 — Conclusion : matrice
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["ccl"])
    bloc_texte(s, "Quelle stratégie pour quel club ?", [
        "Club qui démarre, aucun historique → modèle génératif, sans entraînement.",
        "Club en phase d'installation, 10 à 30 mails par catégorie → SetFit.",
        "Club établi, plus de 100 mails par catégorie → méthode historique ou CamemBERT.",
        "Club exigeant sur la qualité à coût maîtrisé → système hybride, seuil ajustable.",
    ], w=12.1, h=2.6)
    encart(s, "💡 Il n'existe pas de meilleur modèle universel : la bonne réponse dépend du "
              "volume de données disponible et des contraintes de coût.", y=4.3, h=0.9)

    # 27 — Conclusion : limites
    s = nouvelle(prs, L_CONTENU); titre_section(s, SECTIONS["ccl"])
    bloc_texte(s, "Limites et suites envisageables", [
        "Le jeu de données est synthétique : un corpus réel serait probablement plus ambigu.",
        "Le modèle le plus performant à faibles volumes n'est pas celui qui est déployé, "
        "faute d'infrastructure adaptée.",
        "Pistes d'amélioration :",
        ("valider sur des mails réels anonymisés ;",),
        ("héberger SetFit sur un service dédié et l'appeler depuis le tableau de bord ;",),
        ("mesurer la calibration des scores de confiance ;",),
        ("faire étiqueter en priorité les mails les plus incertains.",),
    ], w=12.1)

    # 28 — Merci
    s = nouvelle(prs, L_FIN)
    titre_section(s, "Merci de votre attention")

    prs.save(str(OUT))
    print(f"Support généré : {OUT}")
    print(f"Diapositives : {len(prs.slides.__iter__.__self__._sldIdLst)}")


if __name__ == "__main__":
    build()
